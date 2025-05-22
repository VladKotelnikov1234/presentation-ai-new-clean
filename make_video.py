import os
from pptx import Presentation
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
from mutagen.mp3 import MP3
import logging
from PIL import Image, ImageDraw, ImageFont
import shutil

# Настройка логирования
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s %(levelname)s %(message)s',
    handlers=[
        logging.FileHandler('C:/Users/79824/Downloads/video_creation.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Пути к файлам
PPTX_PATH = 'C:/Users/79824/Downloads/Презентация.pptx'
AUDIO_FILES = [
    'C:/Users/79824/Downloads/slide1.mp3',
    'C:/Users/79824/Downloads/slide2.mp3',
    'C:/Users/79824/Downloads/slide3.mp3',
    'C:/Users/79824/Downloads/slide4.mp3'
]
OUTPUT_VIDEO = 'C:/Users/79824/Downloads/output.mp4'
TEMP_DIR = 'C:/Users/79824/Downloads/temp_slides'
VIDEO_WIDTH, VIDEO_HEIGHT = 1920, 1080


def check_files():
    """Проверяет наличие файлов."""
    if not os.path.exists(PPTX_PATH):
        logger.error(f"Файл презентации не найден: {PPTX_PATH}")
        raise FileNotFoundError(f"Презентация не найдена: {PPTX_PATH}")
    for audio in AUDIO_FILES:
        if not os.path.exists(audio):
            logger.error(f"Аудиофайл не найден: {audio}")
            raise FileNotFoundError(f"Аудиофайл не найден: {audio}")
    logger.info("Все файлы найдены")


def extract_slide_content():
    """Извлекает текст и базовые данные из слайдов."""
    try:
        prs = Presentation(PPTX_PATH)
        slides_content = []
        for slide in prs.slides:
            slide_data = {'text': [], 'background': None}
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data['text'].append(shape.text)
            # Пытаемся извлечь цвет фона (если есть)
            if slide.background.fill.solid():
                rgb = slide.background.fill.fore_color.rgb
                slide_data['background'] = (rgb[0], rgb[1], rgb[2])
            else:
                slide_data['background'] = (255, 255, 255)  # Белый по умолчанию
            slides_content.append(slide_data)
        logger.info(f"Извлечено {len(slides_content)} слайдов")
        return slides_content
    except Exception as e:
        logger.error(f"Ошибка чтения презентации: {e}")
        raise


def render_slide_image(slide_data, slide_number):
    """Рендерит слайд как изображение."""
    try:
        slide_path = os.path.join(TEMP_DIR, f'slide_{slide_number}.png')
        img = Image.new('RGB', (VIDEO_WIDTH, VIDEO_HEIGHT), color=slide_data['background'])
        draw = ImageDraw.Draw(img)

        # Рендерим текст
        try:
            font = ImageFont.truetype("arial.ttf", 60)
        except:
            font = ImageFont.load_default()

        y_offset = 100
        for text in slide_data['text']:
            lines = text.split('\n')
            for line in lines:
                if line.strip():
                    draw.text((100, y_offset), line, fill=(0, 0, 0), font=font)
                    y_offset += 80

        img.save(slide_path, 'PNG')
        logger.info(f"Сохранён слайд {slide_number}: {slide_path}")
        return slide_path
    except Exception as e:
        logger.error(f"Ошибка рендеринга слайда {slide_number}: {e}")
        return None


def get_audio_duration(audio_path):
    """Возвращает длительность аудиофайла."""
    try:
        audio = MP3(audio_path)
        duration = audio.info.length
        logger.info(f"Длительность {audio_path}: {duration:.2f} секунд")
        return duration
    except Exception as e:
        logger.error(f"Ошибка чтения длительности {audio_path}: {e}")
        raise


def create_video():
    """Создаёт видео."""
    # Проверяем файлы
    check_files()

    # Извлекаем содержимое слайдов
    slides_content = extract_slide_content()

    # Проверяем слайды
    if len(slides_content) < 4:
        logger.error(f"Недостаточно слайдов: {len(slides_content)}")
        raise ValueError("Недостаточно слайдов")

    # Создаём временную папку
    os.makedirs(TEMP_DIR, exist_ok=True)

    # Рендерим слайды
    image_paths = []
    for i, slide_data in enumerate(slides_content[:4], 1):
        slide_path = render_slide_image(slide_data, i)
        if slide_path:
            image_paths.append(slide_path)

    # Собираем клипы
    clips = []
    for i, slide_path in enumerate(image_paths):
        if os.path.exists(slide_path):
            duration = get_audio_duration(AUDIO_FILES[i])
            try:
                clip = ImageClip(slide_path, duration=duration).set_audio(AudioFileClip(AUDIO_FILES[i]))
                logger.info(f"Добавлен слайд {i + 1} с {AUDIO_FILES[i]}")
                clips.append(clip)
            except Exception as e:
                logger.error(f"Ошибка создания клипа для слайда {i + 1}: {e}")
                continue

    # Объединяем клипы
    if not clips:
        logger.error("Нет клипов")
        raise ValueError("Нет клипов")

    try:
        logger.info("Объединяю клипы")
        final_clip = concatenate_videoclips(clips, method='compose')
        logger.info(f"Сохраняю видео в {OUTPUT_VIDEO}")
        final_clip.write_videofile(
            OUTPUT_VIDEO,
            codec='libx264',
            audio_codec='aac',
            fps=30,
            threads=4
        )
        logger.info("Видео создано")
    except Exception as e:
        logger.error(f"Ошибка создания видео: {e}")
        raise


def cleanup():
    """Удаляет временные файлы."""
    if os.path.exists(TEMP_DIR):
        shutil.rmtree(TEMP_DIR)
        logger.info("Временные файлы удалены")


if __name__ == '__main__':
    try:
        create_video()
    except Exception as e:
        logger.error(f"Ошибка: {e}")
    finally:
        cleanup()